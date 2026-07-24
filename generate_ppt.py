import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide size to Widescreen 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_BG = RGBColor(15, 23, 42)       # Slate 900
    CARD_BG = RGBColor(30, 41, 59)       # Slate 800
    ACCENT_BLUE = RGBColor(56, 189, 248) # Sky 400
    ACCENT_PURPLE = RGBColor(168, 85, 247) # Purple 500
    TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    WHITE = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(51, 65, 85)   # Slate 700
    GREEN_ACCENT = RGBColor(52, 211, 153) # Emerald 400
    AMBER_ACCENT = RGBColor(251, 191, 36) # Amber 400
    
    blank_slide_layout = prs.slide_layouts[6]

    def add_background(slide, color=DARK_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.color.rgb = color
        return bg

    def add_header(slide, title_text, category_text="ENTERPRISE SRE & IT OPERATIONS INTELLIGENCE PLATFORM"):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.color.rgb = CARD_BORDER

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide1)
    
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "🛡️ ENTERPRISE SRE & IT OPERATIONS INTELLIGENCE"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_BLUE
    
    p1 = tf1.add_paragraph()
    p1.text = "Agent Performance Analyzer"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = "Air-Gapped, Offline-First SRE Dashboard & Local AI Analytics Platform"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    
    add_card(slide1, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.5))
    b_box = slide1.shapes.add_textbox(Inches(1.3), Inches(5.4), Inches(10.733), Inches(1.1))
    btf = b_box.text_frame
    btf.word_wrap = True
    
    bp0 = btf.paragraphs[0]
    bp0.text = "Key Operational Pillars:"
    bp0.font.size = Pt(12)
    bp0.font.bold = True
    bp0.font.color.rgb = ACCENT_PURPLE
    
    bp1 = btf.add_paragraph()
    bp1.text = "• 100% Local SQLite Staging & Ollama Local LLM Inference (Zero Cloud Data Leakage)"
    bp1.font.size = Pt(12)
    bp1.font.color.rgb = TEXT_LIGHT
    
    bp2 = btf.add_paragraph()
    bp2.text = "• Dynamic Context-Aware SLA Performance Scoring (Incident vs. Service Request Weights)"
    bp2.font.size = Pt(12)
    bp2.font.color.rgb = TEXT_LIGHT

    bp3 = btf.add_paragraph()
    bp3.text = "• Automated Infrastructure Noise Clustering & Executive PDF/HTML Review Compiler"
    bp3.font.size = Pt(12)
    bp3.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 2: AGENDA
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide2)
    add_header(slide2, "Executive Presentation Agenda")
    
    agenda_items = [
        ("01", "Executive Summary & Problem Statement", "Understanding the need for air-gapped SRE analytics."),
        ("02", "End-to-End Data Pipeline (Input / Output)", "Ingestion from raw CSVs to SQLite staging and reporting."),
        ("03", "SLA Compliance & Diagnostic Engine", "Auditing resolution times against P0-P3 severity limits."),
        ("04", "Performance Score Math & Dynamic Weights", "Detailed mathematical formulation of SRE matrix scoring."),
        ("05", "Backend Mathematical & SLA Implementation", "Code-level equations for SLA compliance & vector scoring."),
        ("06", "Infrastructure Noise Clustering & Playbooks", "Top 5 systemic alerts and local LLM security playbooks."),
        ("07", "Forensic Ticket Investigator & AI Coaching", "Semantic vector search and localized agent coaching."),
        ("08", "Automated Executive Reporting Engine", "Dynamic HTML/PDF executive review generation."),
        ("09", "System Architecture & Scalability Roadmap", "Tech stack, Streamlit UI, and PostgreSQL/Vector DB roadmap.")
    ]
    
    col_w = Inches(3.7)
    card_h = Inches(1.5)
    
    for idx, (num, title, desc) in enumerate(agenda_items):
        col = idx % 3
        row = idx // 3
        
        left = Inches(0.8) + col * Inches(3.95)
        top = Inches(1.6) + row * Inches(1.7)
        
        add_card(slide2, left, top, col_w, card_h)
        
        tb = slide2.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), col_w - Inches(0.3), card_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{num}. {title}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        p_sub = tf.add_paragraph()
        p_sub.text = desc
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: DATA PIPELINE (INPUT AND OUTPUT)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide3)
    add_header(slide3, "End-to-End Data Pipeline: Inputs & Outputs")
    
    add_card(slide3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3))
    tb_in = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_in = tb_in.text_frame
    tf_in.word_wrap = True
    
    p = tf_in.paragraphs[0]
    p.text = "📥 INPUT DATA SPECIFICATION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    
    inputs_list = [
        "Raw Ticket Exports (data/*.csv):",
        "• Ticket ID & Timestamps (Created Time, Resolved Time)",
        "• Ticket Metadata (Subject, Description, Priority, Status)",
        "• Operational Context (Company, Ticket Type / Category)",
        "• Human Effort Metrics (Effort Required in Mins, Resolution Hours)",
        "• Agent Identity & Resolution Notes",
        "",
        "Data Ingestion Resilience (utils/loader.py):",
        "• Dynamic header parsing & numerical coercion",
        "• Automated timestamp difference fallback if Resolution Hours missing",
        "• Purges staging cache & performs batch upserts into SQLite DB"
    ]
    for line in inputs_list:
        p_item = tf_in.add_paragraph()
        p_item.text = line
        if line.startswith("Raw") or line.startswith("Data"):
            p_item.font.bold = True
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = WHITE
        else:
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = TEXT_LIGHT

    add_card(slide3, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.3))
    tb_out = slide3.shapes.add_textbox(Inches(7.133), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_out = tb_out.text_frame
    tf_out.word_wrap = True
    
    p = tf_out.paragraphs[0]
    p.text = "📤 OUTPUT ARTIFACTS & DELIVERABLES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    
    outputs_list = [
        "Interactive Operations Dashboard (Streamlit UI):",
        "• Real-time reactive Leaderboard & SRE performance ranking",
        "• Dual SLA Tracking Tables (Compliant vs Breached tickets)",
        "• Interactive Plotly distribution & workload charts",
        "",
        "AI Forensic & Security Playbooks:",
        "• Top 5 Infrastructure Noise Alerts & security playbooks",
        "• LLM Ticket Forensic Diagnostics (Root Cause & Summaries)",
        "• Personalized Local AI Career Coaching Profiles",
        "",
        "Executive Exports (reports/):",
        "• Standalone Executive HTML & PDF Operations Review Reports"
    ]
    for line in outputs_list:
        p_item = tf_out.add_paragraph()
        p_item.text = line
        if line.endswith(":"):
            p_item.font.bold = True
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = WHITE
        else:
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 4: SLA COMPLIANCE ENGINE
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide4)
    add_header(slide4, "SLA Compliance Engine & Audit Rules")
    
    priorities = [
        ("URGENT (P0)", "4.0 Hours", "Critical infrastructure outages, high security risks.", ACCENT_PURPLE),
        ("HIGH (P1)", "8.0 Hours", "Core service degradation affecting enterprise workflows.", ACCENT_BLUE),
        ("MEDIUM (P2)", "16.0 Hours", "Standard operational requests and minor service issues.", GREEN_ACCENT),
        ("LOW (P3)", "24.0 Hours", "Informational queries, access requests, general routine tasks.", AMBER_ACCENT)
    ]
    
    card_w = Inches(2.7)
    card_h = Inches(2.2)
    for idx, (p_name, p_target, p_desc, color) in enumerate(priorities):
        left = Inches(0.8) + idx * Inches(2.95)
        add_card(slide4, left, Inches(1.6), card_w, card_h)
        
        tb = slide4.shapes.add_textbox(left + Inches(0.15), Inches(1.75), card_w - Inches(0.3), card_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = p_name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_t = tf.add_paragraph()
        p_t.text = p_target
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        
        p_d = tf.add_paragraph()
        p_d.text = p_desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED

    add_card(slide4, Inches(0.8), Inches(4.1), Inches(11.733), Inches(2.8))
    tb_b = slide4.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.333), Inches(2.4))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    
    p = tf_b.paragraphs[0]
    p.text = "🔍 SLA AUDIT & DIAGNOSTIC LOGIC (analytics/sla.py)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    sla_bullets = [
        "• Timestamp Precision: Compares ticket resolution duration (resolution_hours) strictly against SLA_TARGETS[priority].",
        "• Compliance Binary Mapping: Tags every ticket with sla_breached = 0 (Compliant) or sla_breached = 1 (Breached).",
        "• Global System Audit: Automatically flags missing priority tags, assigning default medium priority rules where unassigned.",
        "• Real-time Metric Aggregation: Computes team-wide compliance percentage: Compliance Rate = (Compliant Tickets / Total Tickets) * 100."
    ]
    for bullet in sla_bullets:
        p_item = tf_b.add_paragraph()
        p_item.text = bullet
        p_item.font.size = Pt(12)
        p_item.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 5: PERFORMANCE SCORE FORMULA & DYNAMIC WEIGHTS
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide5)
    add_header(slide5, "Performance Score Formula & Weight Matrix")
    
    add_card(slide5, Inches(0.8), Inches(1.6), Inches(11.733), Inches(1.3), bg_color=CARD_BG, border_color=ACCENT_BLUE)
    tb_form = slide5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.333), Inches(1.0))
    tf_form = tb_form.text_frame
    tf_form.word_wrap = True
    
    pf0 = tf_form.paragraphs[0]
    pf0.text = "CORE SRE PERFORMANCE SCORE MATHEMATICAL FORMULA"
    pf0.font.size = Pt(11)
    pf0.font.bold = True
    pf0.font.color.rgb = ACCENT_BLUE
    
    pf1 = tf_form.add_paragraph()
    pf1.text = "Score = [ (SLA_Score × w_sla) + (Speed_Score × w_speed) + (Vol_Score × w_vol) + (Effort_Score × w_effort) ] × 100"
    pf1.font.size = Pt(14)
    pf1.font.bold = True
    pf1.font.color.rgb = GREEN_ACCENT

    # 3 Weight Cards
    weight_contexts = [
        ("INCIDENT MODE", "Prioritizes Speed & SLA Limits", [
            ("SLA Compliance (w_sla)", "50%", ACCENT_PURPLE),
            ("Resolution Speed (w_speed)", "35%", ACCENT_BLUE),
            ("Volume Impact (w_vol)", "10%", TEXT_MUTED),
            ("Effort Efficiency (w_effort)", "5%", TEXT_MUTED)
        ], ACCENT_PURPLE),
        
        ("SERVICE REQUEST (SR)", "Prioritizes Volume Throughput", [
            ("Volume Impact (w_vol)", "40%", GREEN_ACCENT),
            ("SLA Compliance (w_sla)", "20%", ACCENT_PURPLE),
            ("Resolution Speed (w_speed)", "20%", ACCENT_BLUE),
            ("Effort Efficiency (w_effort)", "20%", AMBER_ACCENT)
        ], GREEN_ACCENT),
        
        ("ALL TYPES (BALANCED)", "Standard Baseline Matrix", [
            ("SLA Compliance (w_sla)", "40%", ACCENT_PURPLE),
            ("Resolution Speed (w_speed)", "30%", ACCENT_BLUE),
            ("Volume Impact (w_vol)", "20%", GREEN_ACCENT),
            ("Effort Efficiency (w_effort)", "10%", AMBER_ACCENT)
        ], ACCENT_BLUE)
    ]
    
    for idx, (w_title, w_subtitle, w_weights, border_c) in enumerate(weight_contexts):
        left = Inches(0.8) + idx * Inches(3.95)
        top = Inches(3.1)
        card_w = Inches(3.8)
        card_h = Inches(3.8)
        
        add_card(slide5, left, top, card_w, card_h, border_color=border_c)
        tb = slide5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = w_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = border_c
        
        ps = tf.add_paragraph()
        ps.text = w_subtitle
        ps.font.size = Pt(10)
        ps.font.color.rgb = TEXT_MUTED
        
        tf.add_paragraph()
        
        for w_name, w_val, w_color in w_weights:
            pw = tf.add_paragraph()
            pw.text = f"• {w_name}: {w_val}"
            pw.font.size = Pt(11)
            pw.font.bold = True
            pw.font.color.rgb = w_color

    # ==========================================
    # NEW SLIDE 6: BACKEND MATHEMATICAL IMPLEMENTATION DEEP-DIVE
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide6)
    add_header(slide6, "Backend Mathematical & SLA Implementation Deep-Dive")
    
    # Card 1: SLA Backend Engine Logic (Top Half)
    add_card(slide6, Inches(0.8), Inches(1.6), Inches(11.733), Inches(2.6))
    tb_m1 = slide6.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(11.333), Inches(2.3))
    tf_m1 = tb_m1.text_frame
    tf_m1.word_wrap = True
    
    p = tf_m1.paragraphs[0]
    p.text = "⚡ 1. BACKEND SLA AUDIT MATH IMPLEMENTATION (analytics/sla.py & utils/loader.py)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    m1_bullets = [
        "• Timestamp Hour Delta:  Resolution_Hours = max(0.0, (Timestamp_Resolved - Timestamp_Created) / 3600.0)",
        "• Binary SLA Breach Condition:  sla_breached = 1  IF  Resolution_Hours > SLA_TARGETS[Priority]  ELSE  0",
        "• Target Limit Map:  { Urgent (P0): 4.0h, High (P1): 8.0h, Medium (P2): 16.0h, Low (P3): 24.0h }",
        "• Global Team Compliance Pct:  Compliance_Pct = [ (Total_Tickets - sum(sla_breached)) / Total_Tickets ] × 100",
        "• Persistence: Database SQL transaction execute update  'UPDATE tickets SET sla_breached = ? WHERE ticket_id = ?'"
    ]
    for b in m1_bullets:
        p_item = tf_m1.add_paragraph()
        p_item.text = b
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = TEXT_LIGHT

    # Card 2: Leaderboard Normalization & Scoring Math (Bottom Half)
    add_card(slide6, Inches(0.8), Inches(4.4), Inches(11.733), Inches(2.6))
    tb_m2 = slide6.shapes.add_textbox(Inches(1.0), Inches(4.55), Inches(11.333), Inches(2.3))
    tf_m2 = tb_m2.text_frame
    tf_m2.word_wrap = True
    
    p2 = tf_m2.paragraphs[0]
    p2.text = "📊 2. VECTOR NORMALIZATION & MATRIX SCORING IMPLEMENTATION (analytics/scoring.py)"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = GREEN_ACCENT
    
    m2_bullets = [
        "• Pandas Groupby Aggregation:  summary = df.groupby('agent').agg(Tickets_Handled=('ticket_id','count'), ...)",
        "• Linear Relative Normalization:  vol_score = Tickets_Handled / max(Tickets_Handled)",
        "• Inverted Metric Scaling (Speed & Effort):  speed_score = 1 - (Avg_Resolution_Hours / max(Avg_Resolution_Hours))",
        "                                              effort_score = 1 - (Avg_Effort_Mins / max(Avg_Effort_Mins))",
        "• Weighted Dot Product Execution:  Performance_Score = ( (sla_score × w_sla) + (speed_score × w_speed) + ",
        "                                                       (vol_score × w_vol) + (effort_score × w_effort) ) × 100"
    ]
    for b in m2_bullets:
        p_item = tf_m2.add_paragraph()
        p_item.text = b
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 7: INFRASTRUCTURE NOISE CLUSTERING
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide7)
    add_header(slide7, "Infrastructure Noise Clustering & AI Playbooks")
    
    add_card(slide7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3))
    tb_l = slide7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "🚨 TOP 5 SYSTEMIC ALERTS (analytics/noise_cluster.py)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    
    noise_bullets = [
        "Systemic Pattern Frequency Mining:",
        "• Aggregates high-frequency repetitive ticket subjects across target companies.",
        "• Surfaces Top 5 highest frequency operational alerts (e.g. repeated VPN disconnects, password lockouts, storage quotas).",
        "",
        "Operational Value:",
        "• Prevents team burnout by isolating recurring infrastructure bugs.",
        "• Distinguishes true operational incidents from background noise.",
        "• Provides executive visibility into enterprise company risk profiles."
    ]
    for bullet in noise_bullets:
        p_item = tf_l.add_paragraph()
        p_item.text = bullet
        if bullet.endswith(":"):
            p_item.font.bold = True
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = WHITE
        else:
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = TEXT_LIGHT

    add_card(slide7, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.3))
    tb_r = slide7.shapes.add_textbox(Inches(7.133), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "🧠 AIR-GAPPED SECURITY PLAYBOOKS (SystemicRootCauseEngine)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    
    ai_bullets = [
        "Local LLM Inference (qwen2.5:3b):",
        "• Transmits noise cluster text payloads strictly to local Ollama API endpoint.",
        "• Evaluates security posture impact and engineering effort drain.",
        "",
        "Automated Playbook Generation:",
        "• Constructs actionable mitigation playbooks directly inside the UI.",
        "• Provides automated root-cause summaries without cloud LLM subscriptions.",
        "• Guarantees complete data privacy compliance for air-gapped deployments."
    ]
    for bullet in ai_bullets:
        p_item = tf_r.add_paragraph()
        p_item.text = bullet
        if bullet.endswith(":"):
            p_item.font.bold = True
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = WHITE
        else:
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 8: FORENSIC INVESTIGATOR & COACHING
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide8)
    add_header(slide8, "Forensic Ticket Investigator & AI Coaching")
    
    add_card(slide8, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3))
    tb_f = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    
    p = tf_f.paragraphs[0]
    p.text = "🔍 AI FORENSIC AUDITOR (analytics/ticket_explorer.py)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    forensic_bullets = [
        "Semantic Vector Search (LocalTicketVectorStore):",
        "• Computes TF-IDF / Cosine similarity vectors over ticket corpus natively.",
        "• Instantly surfaces top historically solved similar cases for SREs.",
        "",
        "AI Forensic Diagnostics (LocalTicketAnalyzer):",
        "• Analyzes raw descriptions and resolution notes via local LLM.",
        "• Outputs structured root-cause diagnostic summaries.",
        "• Eliminates duplicate engineering effort on recurring issues."
    ]
    for bullet in forensic_bullets:
        p_item = tf_f.add_paragraph()
        p_item.text = bullet
        if bullet.endswith(":"):
            p_item.font.bold = True
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = WHITE
        else:
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = TEXT_LIGHT

    add_card(slide8, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.3))
    tb_c = slide8.shapes.add_textbox(Inches(7.133), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    
    p = tf_c.paragraphs[0]
    p.text = "🎓 CAREER COACHING WORKSHOP (LocalAgentCoachingEngine)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    
    coach_bullets = [
        "Individual SRE Profile Audits:",
        "• Batches an individual engineer's ticket history into LLM context window.",
        "• Evaluates SLA breach history, effort allocation, and resolution notes quality.",
        "",
        "Actionable Growth Feedback:",
        "• Generates tailored coaching remarks highlighting operational strengths.",
        "• Recommends specific technical areas for SRE skill development.",
        "• Empowers engineering leadership with automated performance reviews."
    ]
    for bullet in coach_bullets:
        p_item = tf_c.add_paragraph()
        p_item.text = bullet
        if bullet.endswith(":"):
            p_item.font.bold = True
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = WHITE
        else:
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 9: SYSTEM ARCHITECTURE & TECH STACK
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide9)
    add_header(slide9, "System Architecture & Technical Stack")
    
    tech_stack = [
        ("Frontend & UI Engine", "Streamlit", "Reactive UI with dual-theme glassmorphism, responsive metric cards, and CSS menu overrides.", ACCENT_BLUE),
        ("Data & Storage Layer", "Pandas + SQLite3", "Localized fast-read SQLite staging gateway (analyzer.db) with pre-compiled B-tree indexes.", GREEN_ACCENT),
        ("Air-Gapped AI Inference", "Ollama (qwen2.5:3b)", "Self-hosted local LLM endpoint running at localhost:11434 with strict JSON schema parsing.", ACCENT_PURPLE),
        ("Document Export Engine", "xhtml2pdf + HTML5", "Enterprise SLA report generator compiling HTML templates into paginated executive PDFs.", AMBER_ACCENT)
    ]
    
    for idx, (cat_title, tech_name, tech_desc, color) in enumerate(tech_stack):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8) + col * Inches(5.95)
        top = Inches(1.6) + row * Inches(2.7)
        
        add_card(slide9, left, top, Inches(5.75), Inches(2.4))
        tb = slide9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.35), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = cat_title.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        
        pt = tf.add_paragraph()
        pt.text = tech_name
        pt.font.size = Pt(20)
        pt.font.bold = True
        pt.font.color.rgb = color
        
        pd = tf.add_paragraph()
        pd.text = tech_desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 10: FUTURE ROADMAP & SCALABILITY
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_background(slide10)
    add_header(slide10, "Future Scalability & Engineering Roadmap")
    
    roadmap_items = [
        ("1. PostgreSQL Migration", "Transition from SQLite to PostgreSQL", "Enables high-throughput concurrent writes from automated webhooks while maintaining SQL compatibility via SQLAlchemy.", ACCENT_PURPLE),
        ("2. Asynchronous LLM Pipeline", "Migrate to asyncio & httpx", "Decouples local Ollama LLM inference from Streamlit main UI looper thread, preventing UI blocking on complex prompts.", ACCENT_BLUE),
        ("3. Dedicated Vector Database", "Integrate ChromaDB / Qdrant", "Replaces native similarity calculations with high-density vector database index for 100k+ ticket semantic search.", GREEN_ACCENT),
        ("4. Distributed Cache Layer", "Implement Redis Caching", "Replaces st.cache_data with distributed Redis cache to support multi-replica containerized deployments.", AMBER_ACCENT)
    ]
    
    card_w = Inches(5.6)
    card_h = Inches(2.4)
    
    for idx, (num_title, main_title, desc, color) in enumerate(roadmap_items):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8) if col == 0 else Inches(6.933)
        top = Inches(1.6) + row * Inches(2.7)
        
        add_card(slide10, left, top, card_w, card_h)
        tb = slide10.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), card_w - Inches(0.4), card_h - Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = num_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        pt = tf.add_paragraph()
        pt.text = main_title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = WHITE
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_LIGHT

    # Save presentation
    output_path = "/home/rauf/agent-performance-analyzer/Enterprise_SRE_Operations_Platform.pptx"
    prs.save(output_path)
    print(f"✅ Presentation successfully generated and saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
